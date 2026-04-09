from pathlib import Path

from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.table import WD_ALIGN_VERTICAL, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Mm, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt


ROOT = Path("/Users/shimane/Projects/Plug-A-Pro/docs/kgolaentle")
ASSETS = ROOT / "assets"
OUTPUT = ROOT / "output"
OUTPUT.mkdir(parents=True, exist_ok=True)

PRIMARY = "9B1B30"
PRIMARY_DARK = "7A1526"
SECONDARY = "0F0F1A"
ACCENT = "D4A843"
ACCENT_LIGHT = "E8C876"
WARM_GRAY = "F7F5F2"
WARM_WHITE = "FDFCFB"
BODY = "333333"
MUTED = "6A7282"

DOCX_PATH = OUTPUT / "Kgolaentle Holdings Company Profile - Procurement Edition.docx"
PPTX_PATH = OUTPUT / "Kgolaentle Holdings Company Profile - Procurement Edition.pptx"

LOGO = ASSETS / "logo.png"
HERO = ASSETS / "hero-serving-lunch.jpg"

PORTFOLIOS = [
    {
        "title": "Kgolaentle Rentals",
        "tagline": "Premium event and practical rental capability",
        "image": ASSETS / "service-rentals.jpg",
        "summary": (
            "Kgolaentle Rentals supports events, functions and practical site needs with "
            "dependable rental solutions built around comfort, timing and service discipline."
        ),
        "serves": (
            "Private clients, event organisers, community functions and corporate gatherings "
            "that need reliable equipment and a partner that manages the details properly."
        ),
        "value": (
            "Clients gain confidence that the equipment will arrive ready, the setup will be "
            "handled well and the experience will reflect positively on the event."
        ),
        "bullets": [
            "VIP mobile toilets for events and functions",
            "Mobile freezer rentals",
            "Packages shaped around budget and operational need",
            "Professional setup and maintenance support",
        ],
    },
    {
        "title": "Courier Services",
        "tagline": "Trusted regional delivery with local knowledge",
        "image": ASSETS / "service-courier.jpg",
        "summary": (
            "Courier Services provides dependable delivery capability across Rustenburg and "
            "surrounding areas through strong route knowledge, communication discipline and "
            "consistent execution."
        ),
        "serves": (
            "Individuals, households and businesses that need a regional delivery partner close "
            "to the communities they operate in."
        ),
        "value": (
            "The portfolio delivers confidence in movement. Parcels are handled with care, "
            "communication stays clear and service remains rooted in local context."
        ),
        "bullets": [
            "Coverage across Rustenburg and surrounding areas",
            "Deep local knowledge across Sun City, Ledig, Mogwase and nearby communities",
            "Online tracking and timely communication",
            "Professionally trained delivery team",
        ],
    },
    {
        "title": "Technology Solutions",
        "tagline": "Strategy, systems and execution for real-world technology work",
        "image": ASSETS / "service-technology.jpg",
        "summary": (
            "Technology Solutions helps founders, operators and enterprise teams design, build "
            "and deploy digital platforms that work in the real world."
        ),
        "serves": (
            "Founders, operators and enterprise teams that need product strategy, technical "
            "leadership and disciplined execution on complex digital work."
        ),
        "value": (
            "Clients gain a partner that can move from strategy to architecture to working "
            "delivery without losing business context or execution quality."
        ),
        "bullets": [
            "Digital platform design and product strategy",
            "Web and mobile application delivery",
            "AI-enabled workflows and automation",
            "Enterprise architecture and technology leadership",
        ],
    },
    {
        "title": "Opulent Beauty",
        "tagline": "Premium beauty experiences delivered with care",
        "image": ASSETS / "service-beauty.jpg",
        "summary": (
            "Opulent Beauty brings together premium services and curated products in an "
            "environment shaped by professionalism, care and attention to client experience."
        ),
        "serves": (
            "Clients seeking a refined beauty experience as well as trusted products and "
            "treatments delivered in a welcoming professional setting."
        ),
        "value": (
            "The portfolio creates confidence through thoughtful service, polished presentation "
            "and a standard that feels intentional from start to finish."
        ),
        "bullets": [
            "Premium beauty services and treatments",
            "Curated beauty products",
            "Professional and welcoming environment",
            "Client experience shaped around quality and care",
        ],
    },
]


def rgb(hex_code: str) -> RGBColor:
    return RGBColor.from_string(hex_code)


def ppt_rgb(hex_code: str) -> PptRGBColor:
    return PptRGBColor.from_string(hex_code)


def set_cell_shading(cell, fill):
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), fill)
    tc_pr.append(shd)


def set_cell_border(cell, **kwargs):
    tc = cell._tc
    tc_pr = tc.get_or_add_tcPr()
    tc_borders = tc_pr.first_child_found_in("w:tcBorders")
    if tc_borders is None:
        tc_borders = OxmlElement("w:tcBorders")
        tc_pr.append(tc_borders)
    for edge in ("left", "top", "right", "bottom"):
        edge_data = kwargs.get(edge)
        if edge_data:
            tag = f"w:{edge}"
            element = tc_borders.find(qn(tag))
            if element is None:
                element = OxmlElement(tag)
                tc_borders.append(element)
            for key, value in edge_data.items():
                element.set(qn(f"w:{key}"), str(value))


def remove_table_borders(table):
    tbl = table._tbl
    tbl_pr = tbl.tblPr
    borders = tbl_pr.first_child_found_in("w:tblBorders")
    if borders is None:
        borders = OxmlElement("w:tblBorders")
        tbl_pr.append(borders)
    for edge in ("top", "left", "bottom", "right", "insideH", "insideV"):
        element = borders.find(qn(f"w:{edge}"))
        if element is None:
            element = OxmlElement(f"w:{edge}")
            borders.append(element)
        element.set(qn("w:val"), "nil")


def add_paragraph(
    container,
    text,
    *,
    style=None,
    size=11,
    color=BODY,
    bold=False,
    italic=False,
    font="Arial",
    align=WD_ALIGN_PARAGRAPH.LEFT,
    space_after=6,
    space_before=0,
    line_spacing=1.25,
):
    paragraph = container.add_paragraph(style=style)
    paragraph.alignment = align
    paragraph.paragraph_format.space_after = Pt(space_after)
    paragraph.paragraph_format.space_before = Pt(space_before)
    paragraph.paragraph_format.line_spacing = line_spacing
    run = paragraph.add_run(text)
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = rgb(color)
    run.font.bold = bold
    run.font.italic = italic
    return paragraph


def add_footer(section):
    footer = section.footer
    p = footer.paragraphs[0]
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("Kgolaentle Holdings | +27 (0) 87 093 7316 | info@kgolaentle.com | www.kgolaentle.com")
    run.font.name = "Arial"
    run.font.size = Pt(8.5)
    run.font.color.rgb = rgb(MUTED)


def add_section_banner(doc, title, subtitle=None):
    table = doc.add_table(rows=1, cols=1)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    remove_table_borders(table)
    cell = table.cell(0, 0)
    cell.vertical_alignment = WD_ALIGN_VERTICAL.CENTER
    set_cell_shading(cell, SECONDARY)
    set_cell_border(
        cell,
        top={"val": "single", "sz": 0, "color": SECONDARY},
        bottom={"val": "single", "sz": 0, "color": SECONDARY},
    )
    p = cell.paragraphs[0]
    p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    p.paragraph_format.space_before = Pt(6)
    p.paragraph_format.space_after = Pt(2)
    r = p.add_run(title)
    r.font.name = "Georgia"
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = rgb(WARM_WHITE)
    if subtitle:
        p2 = cell.add_paragraph()
        p2.alignment = WD_ALIGN_PARAGRAPH.LEFT
        p2.paragraph_format.space_before = Pt(0)
        p2.paragraph_format.space_after = Pt(8)
        r2 = p2.add_run(subtitle)
        r2.font.name = "Arial"
        r2.font.size = Pt(9)
        r2.font.color.rgb = rgb(ACCENT_LIGHT)
    doc.add_paragraph()


def build_docx():
    doc = Document()
    section = doc.sections[0]
    section.page_width = Mm(210)
    section.page_height = Mm(297)
    section.top_margin = Mm(14)
    section.bottom_margin = Mm(14)
    section.left_margin = Mm(16)
    section.right_margin = Mm(16)
    add_footer(section)

    cover = doc.add_table(rows=1, cols=1)
    cover.alignment = WD_TABLE_ALIGNMENT.CENTER
    remove_table_borders(cover)
    cell = cover.cell(0, 0)
    set_cell_shading(cell, SECONDARY)
    p_logo = cell.paragraphs[0]
    p_logo.alignment = WD_ALIGN_PARAGRAPH.LEFT
    run_logo = p_logo.add_run()
    run_logo.add_picture(str(LOGO), width=Inches(1.5))
    add_paragraph(
        cell,
        "Kgolaentle Holdings",
        font="Georgia",
        size=24,
        color=WARM_WHITE,
        bold=True,
        space_after=3,
    )
    add_paragraph(
        cell,
        "Company Profile | Procurement Edition",
        font="Arial",
        size=12,
        color=ACCENT_LIGHT,
        bold=True,
        space_after=2,
    )
    add_paragraph(
        cell,
        "A founder-led group building practical businesses with purpose, discipline and care.",
        font="Arial",
        size=11,
        color=WARM_WHITE,
        space_after=4,
    )
    add_paragraph(
        cell,
        "Aligned to the live Kgolaentle Holdings business structure as of 7 April 2026.",
        font="Arial",
        size=9.5,
        color="D9D9D9",
        space_after=10,
    )

    hero_p = doc.add_paragraph()
    hero_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    hero_run = hero_p.add_run()
    hero_run.add_picture(str(HERO), width=Inches(6.9))

    stats = doc.add_table(rows=1, cols=3)
    stats.alignment = WD_TABLE_ALIGNMENT.CENTER
    remove_table_borders(stats)
    stat_values = [
        ("4", "Distinct portfolios"),
        ("10+", "Years of operating experience"),
        ("North West", "Community-rooted base"),
    ]
    for idx, (value, label) in enumerate(stat_values):
        stat_cell = stats.cell(0, idx)
        stat_cell.vertical_alignment = WD_ALIGN_VERTICAL.CENTER
        set_cell_shading(stat_cell, WARM_GRAY)
        set_cell_border(
            stat_cell,
            top={"val": "single", "sz": 14, "color": ACCENT},
            bottom={"val": "single", "sz": 0, "color": WARM_GRAY},
            left={"val": "single", "sz": 0, "color": WARM_GRAY},
            right={"val": "single", "sz": 0, "color": WARM_GRAY},
        )
        add_paragraph(
            stat_cell,
            value,
            font="Georgia",
            size=20,
            color=PRIMARY,
            bold=True,
            align=WD_ALIGN_PARAGRAPH.CENTER,
            space_after=0,
            space_before=3,
        )
        add_paragraph(
            stat_cell,
            label,
            font="Arial",
            size=9,
            color=MUTED,
            align=WD_ALIGN_PARAGRAPH.CENTER,
            space_after=4,
        )

    doc.add_paragraph()
    overview_note = doc.add_table(rows=1, cols=1)
    remove_table_borders(overview_note)
    cell = overview_note.cell(0, 0)
    set_cell_shading(cell, WARM_WHITE)
    set_cell_border(cell, left={"val": "single", "sz": 18, "color": PRIMARY})
    add_paragraph(
        cell,
        "Procurement view: Kgolaentle Holdings is a diversified group with four clear operating portfolios under one founder-led vision. Each portfolio serves a distinct market while sharing a common standard of trust, execution and long-term discipline.",
        font="Arial",
        size=10.5,
        color=BODY,
        space_before=4,
        space_after=4,
    )

    doc.add_page_break()

    add_section_banner(doc, "Group Overview", "A clear view of the business model and operating direction")
    add_paragraph(
        doc,
        "Kgolaentle Holdings is a diversified South African business group rooted in the North West and built around practical services that improve everyday life. The group brings together Kgolaentle Rentals, Courier Services, Technology Solutions and Opulent Beauty under one leadership vision.",
        size=11,
        space_after=8,
    )
    add_paragraph(
        doc,
        "The business is not a loose collection of services. It is a disciplined holdings model shaped by purpose before products, trust in delivery and growth built for the long term.",
        size=11,
        space_after=10,
    )

    snapshot = doc.add_table(rows=2, cols=2)
    snapshot.alignment = WD_TABLE_ALIGNMENT.CENTER
    remove_table_borders(snapshot)
    for i, portfolio in enumerate(PORTFOLIOS):
        cell = snapshot.cell(i // 2, i % 2)
        set_cell_shading(cell, WARM_GRAY)
        set_cell_border(cell, top={"val": "single", "sz": 14, "color": PRIMARY})
        add_paragraph(
            cell,
            portfolio["title"],
            font="Georgia",
            size=13,
            color=SECONDARY,
            bold=True,
            space_after=2,
            space_before=3,
        )
        add_paragraph(
            cell,
            portfolio["tagline"],
            font="Arial",
            size=9.5,
            color=PRIMARY,
            bold=True,
            space_after=4,
        )
        add_paragraph(
            cell,
            portfolio["summary"],
            font="Arial",
            size=9.5,
            color=BODY,
            space_after=4,
        )

    doc.add_paragraph()
    add_paragraph(
        doc,
        "What procurement teams should see clearly",
        font="Georgia",
        size=15,
        color=SECONDARY,
        bold=True,
        space_after=4,
    )
    for bullet in [
        "Founder-led accountability with a clear group narrative",
        "Four distinct portfolios with defined customer and value positions",
        "Community-rooted operating model with commercial maturity",
        "Service language grounded in execution rather than hype",
    ]:
        p = doc.add_paragraph(style=None)
        p.paragraph_format.left_indent = Inches(0.2)
        p.paragraph_format.space_after = Pt(3)
        run = p.add_run(f"• {bullet}")
        run.font.name = "Arial"
        run.font.size = Pt(10)
        run.font.color.rgb = rgb(BODY)

    doc.add_page_break()

    add_section_banner(doc, "Our Why, Story and Philosophy", "The thinking that holds the group together")
    for title, text in [
        (
            "Our why",
            "Kgolaentle Holdings exists because communities deserve quality, professionalism and dignity in the services they rely on every day.",
        ),
        (
            "Our story",
            "The group began with a clear observation from founder and CEO Masego Mafoko: communities in and around Rustenburg deserved the same standard of service found anywhere else in the country. That belief became the starting point for building a business that grows with discipline and stays close to real needs.",
        ),
        (
            "Our philosophy",
            "Strong businesses are built from the inside out. Purpose comes first. Trust is earned in delivery. Growth must be disciplined. Excellence is the result of consistent standards applied over time.",
        ),
    ]:
        box = doc.add_table(rows=1, cols=1)
        remove_table_borders(box)
        cell = box.cell(0, 0)
        set_cell_shading(cell, WARM_WHITE)
        set_cell_border(cell, left={"val": "single", "sz": 14, "color": ACCENT})
        add_paragraph(
            cell,
            title,
            font="Georgia",
            size=13,
            color=PRIMARY,
            bold=True,
            space_after=2,
            space_before=3,
        )
        add_paragraph(
            cell,
            text,
            font="Arial",
            size=10.5,
            color=BODY,
            space_after=4,
        )
        doc.add_paragraph()

    doc.add_page_break()

    add_section_banner(doc, "Portfolio Detail", "Distinct capabilities under one shared standard")
    for index, portfolio in enumerate(PORTFOLIOS):
        block = doc.add_table(rows=1, cols=2)
        block.alignment = WD_TABLE_ALIGNMENT.CENTER
        remove_table_borders(block)
        img_cell = block.cell(0, 0)
        txt_cell = block.cell(0, 1)
        img_p = img_cell.paragraphs[0]
        img_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        img_p.add_run().add_picture(str(portfolio["image"]), width=Inches(2.55))
        set_cell_shading(txt_cell, WARM_GRAY)
        set_cell_border(txt_cell, top={"val": "single", "sz": 16, "color": PRIMARY})
        add_paragraph(
            txt_cell,
            portfolio["title"],
            font="Georgia",
            size=15,
            color=SECONDARY,
            bold=True,
            space_after=2,
            space_before=4,
        )
        add_paragraph(
            txt_cell,
            portfolio["tagline"],
            font="Arial",
            size=9.5,
            color=PRIMARY,
            bold=True,
            space_after=4,
        )
        add_paragraph(txt_cell, portfolio["summary"], size=9.8, space_after=4)
        add_paragraph(txt_cell, f"Who it serves: {portfolio['serves']}", size=9.5, color=BODY, space_after=3)
        add_paragraph(txt_cell, f"Value delivered: {portfolio['value']}", size=9.5, color=BODY, space_after=3)
        for bullet in portfolio["bullets"]:
            p = txt_cell.add_paragraph()
            p.paragraph_format.left_indent = Inches(0.1)
            p.paragraph_format.space_after = Pt(2)
            r = p.add_run(f"• {bullet}")
            r.font.name = "Arial"
            r.font.size = Pt(9)
            r.font.color.rgb = rgb(BODY)
        doc.add_paragraph()
        if index == 1:
            doc.add_page_break()

    doc.add_page_break()

    add_section_banner(doc, "What Sets The Group Apart", "Procurement-facing strengths and founder-led direction")
    strengths = doc.add_table(rows=1, cols=2)
    remove_table_borders(strengths)
    left = strengths.cell(0, 0)
    right = strengths.cell(0, 1)
    set_cell_shading(left, WARM_GRAY)
    set_cell_shading(right, WARM_WHITE)
    set_cell_border(left, top={"val": "single", "sz": 14, "color": PRIMARY})
    set_cell_border(right, top={"val": "single", "sz": 14, "color": ACCENT})
    add_paragraph(left, "Differentiators", font="Georgia", size=14, color=SECONDARY, bold=True, space_after=4, space_before=4)
    for bullet in [
        "Disciplined diversification rather than opportunistic expansion",
        "Community-rooted operating perspective",
        "Clear standards around service, communication and execution",
        "Commercial maturity without inflated language or empty claims",
    ]:
        p = left.add_paragraph()
        p.paragraph_format.space_after = Pt(3)
        r = p.add_run(f"• {bullet}")
        r.font.name = "Arial"
        r.font.size = Pt(9.5)
        r.font.color.rgb = rgb(BODY)
    add_paragraph(right, "Leadership", font="Georgia", size=14, color=SECONDARY, bold=True, space_after=4, space_before=4)
    add_paragraph(
        right,
        "Kgolaentle Holdings is led by founder and CEO Masego Mafoko. With over ten years of experience in communications, public relations and marketing, she has shaped the group around a simple standard: communities deserve service that is professional, reliable and dignified.",
        size=10,
        space_after=4,
    )
    add_paragraph(
        right,
        "Her leadership has turned the business into more than a collection of ventures. It has created a group with a clear why, a disciplined view of growth and a commitment to building something that lasts.",
        size=10,
        space_after=4,
    )

    doc.add_paragraph()
    add_paragraph(doc, "Our values", font="Georgia", size=15, color=SECONDARY, bold=True, space_after=4)
    values = doc.add_table(rows=1, cols=3)
    remove_table_borders(values)
    for idx, (title, text) in enumerate(
        [
            ("Trust", "We build relationships through reliability, transparency and follow-through."),
            ("Collaboration", "We work closely with clients, partners and teams to achieve better outcomes."),
            ("Continuous improvement", "We keep refining our work, systems and service experience."),
        ]
    ):
        cell = values.cell(0, idx)
        set_cell_shading(cell, WARM_GRAY if idx != 1 else WARM_WHITE)
        set_cell_border(cell, top={"val": "single", "sz": 12, "color": ACCENT if idx == 1 else PRIMARY})
        add_paragraph(cell, title, font="Georgia", size=12, color=PRIMARY, bold=True, space_after=2, space_before=3)
        add_paragraph(cell, text, font="Arial", size=9.2, color=BODY, space_after=3)

    doc.add_page_break()

    add_section_banner(doc, "Closing Statement", "Practical businesses built for long-term trust")
    add_paragraph(
        doc,
        "Kgolaentle Holdings is building a trusted ecosystem of practical businesses for the long term. The ambition is clear: grow with discipline, serve with care and be known for work that people can rely on.",
        size=11,
        space_after=10,
    )

    closing = doc.add_table(rows=1, cols=2)
    remove_table_borders(closing)
    left = closing.cell(0, 0)
    right = closing.cell(0, 1)
    set_cell_shading(left, SECONDARY)
    set_cell_shading(right, WARM_GRAY)
    logo_p = left.paragraphs[0]
    logo_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    logo_p.add_run().add_picture(str(LOGO), width=Inches(1.6))
    add_paragraph(
        left,
        "Excellence rooted in community.",
        font="Georgia",
        size=14,
        color=WARM_WHITE,
        bold=True,
        align=WD_ALIGN_PARAGRAPH.CENTER,
        space_after=4,
    )
    add_paragraph(
        left,
        "Prepared as a procurement-facing company profile aligned to the current live business structure.",
        font="Arial",
        size=9,
        color="D9D9D9",
        align=WD_ALIGN_PARAGRAPH.CENTER,
        space_after=4,
    )
    add_paragraph(right, "Contact details", font="Georgia", size=14, color=SECONDARY, bold=True, space_after=4, space_before=4)
    for line in [
        "Kgolaentle Holdings",
        "Blairgowrie Section, Chaneng",
        "North West, 0310, South Africa",
        "+27 (0) 87 093 7316",
        "info@kgolaentle.com",
        "www.kgolaentle.com",
    ]:
        add_paragraph(right, line, font="Arial", size=10, color=BODY, space_after=2)

    doc.save(DOCX_PATH)


def add_slide_base(slide, *, dark=False):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = ppt_rgb(SECONDARY if dark else WARM_WHITE)
    if dark:
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            PptInches(9.9),
            PptInches(-1.1),
            PptInches(4.8),
            PptInches(4.8),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = ppt_rgb(PRIMARY)
        shape.fill.transparency = 0.55
        shape.line.fill.background()
    else:
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            PptInches(-0.8),
            PptInches(5.5),
            PptInches(3.2),
            PptInches(3.2),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = ppt_rgb(ACCENT_LIGHT)
        shape.fill.transparency = 0.65
        shape.line.fill.background()
    footer = slide.shapes.add_textbox(PptInches(0.55), PptInches(7.08), PptInches(12.2), PptInches(0.22))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "Kgolaentle Holdings | Procurement Edition"
    run.font.name = "Arial"
    run.font.size = PptPt(9)
    run.font.color.rgb = ppt_rgb("D6D6D6" if dark else MUTED)


def add_ppt_title(slide, title, subtitle=None, *, dark=False, left=0.7, top=0.55, width=6.6):
    box = slide.shapes.add_textbox(PptInches(left), PptInches(top), PptInches(width), PptInches(1.4))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    r.font.name = "Georgia"
    r.font.bold = True
    r.font.size = PptPt(24)
    r.font.color.rgb = ppt_rgb(WARM_WHITE if dark else SECONDARY)
    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = "Arial"
        r2.font.size = PptPt(11)
        r2.font.color.rgb = ppt_rgb(ACCENT_LIGHT if dark else PRIMARY)


def add_ppt_body(slide, text, *, left, top, width, height=1.0, size=12, color=BODY, dark=False):
    box = slide.shapes.add_textbox(PptInches(left), PptInches(top), PptInches(width), PptInches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.name = "Arial"
    r.font.size = PptPt(size)
    r.font.color.rgb = ppt_rgb("E9E9E9" if dark else color)
    return tf


def add_ppt_bullets(slide, bullets, *, left, top, width, height, dark=False, size=11):
    box = slide.shapes.add_textbox(PptInches(left), PptInches(top), PptInches(width), PptInches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.bullet = True
        p.space_after = PptPt(4)
        p.font.name = "Arial"
        p.font.size = PptPt(size)
        p.font.color.rgb = ppt_rgb("ECECEC" if dark else BODY)


def add_card(slide, x, y, w, h, title, body, *, accent_color=PRIMARY, dark=False):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ppt_rgb("1B1B28" if dark else WARM_GRAY)
    shape.line.color.rgb = ppt_rgb(accent_color)
    shape.line.width = PptPt(1.6)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Georgia"
    r.font.bold = True
    r.font.size = PptPt(16)
    r.font.color.rgb = ppt_rgb(WARM_WHITE if dark else SECONDARY)
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = body
    r2.font.name = "Arial"
    r2.font.size = PptPt(10.5)
    r2.font.color.rgb = ppt_rgb("EAEAEA" if dark else BODY)


def build_pptx():
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1 cover
    slide = prs.slides.add_slide(blank)
    add_slide_base(slide, dark=True)
    slide.shapes.add_picture(str(HERO), 0, 0, width=prs.slide_width, height=prs.slide_height)
    overlay = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = ppt_rgb(SECONDARY)
    overlay.fill.transparency = 0.28
    overlay.line.fill.background()
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, PptInches(0.7), PptInches(0.85), PptInches(0.12), PptInches(1.6))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ppt_rgb(ACCENT)
    accent.line.fill.background()
    slide.shapes.add_picture(str(LOGO), PptInches(0.95), PptInches(0.5), width=PptInches(1.7))
    add_ppt_title(
        slide,
        "Kgolaentle Holdings",
        "Company Profile | Procurement Edition",
        dark=True,
        left=0.95,
        top=1.55,
        width=6.2,
    )
    add_ppt_body(
        slide,
        "A founder-led group building practical businesses with purpose, discipline and care.",
        left=0.95,
        top=2.55,
        width=5.6,
        height=0.8,
        size=13,
        dark=True,
    )
    add_ppt_body(
        slide,
        "Aligned to the live Kgolaentle Holdings business structure as of 7 April 2026.",
        left=0.95,
        top=3.18,
        width=5.4,
        height=0.5,
        size=10.5,
        color="DDDDDD",
        dark=True,
    )
    for idx, (value, label) in enumerate([("4", "Portfolios"), ("10+", "Years"), ("North West", "Rooted")]):
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PptInches(0.95 + idx * 1.55), PptInches(4.05), PptInches(1.35), PptInches(1.0))
        card.fill.solid()
        card.fill.fore_color.rgb = ppt_rgb("1D1D2A")
        card.fill.transparency = 0.12
        card.line.color.rgb = ppt_rgb(ACCENT)
        tf = card.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = value
        r.font.name = "Georgia"
        r.font.bold = True
        r.font.size = PptPt(16)
        r.font.color.rgb = ppt_rgb(WARM_WHITE)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = label
        r2.font.name = "Arial"
        r2.font.size = PptPt(8.5)
        r2.font.color.rgb = ppt_rgb("EFE3BE")

    # Slide 2 at a glance
    slide = prs.slides.add_slide(blank)
    add_slide_base(slide, dark=False)
    add_ppt_title(slide, "Group At A Glance", "Four distinct portfolios under one founder-led vision")
    add_ppt_body(
        slide,
        "Kgolaentle Holdings is a diversified business group rooted in community and shaped by practical execution. The group does not grow by collecting random services. It grows by building useful businesses that clients can trust.",
        left=0.7,
        top=1.45,
        width=7.1,
        height=0.9,
        size=12,
    )
    for i, portfolio in enumerate(PORTFOLIOS):
        x = 0.7 + (i % 2) * 3.4
        y = 2.45 + (i // 2) * 2.1
        add_card(slide, x, y, 3.05, 1.65, portfolio["title"], portfolio["tagline"], accent_color=PRIMARY if i % 2 == 0 else ACCENT)
    add_card(
        slide,
        7.65,
        1.45,
        4.9,
        4.75,
        "Procurement-facing reading",
        "Founder-led accountability.\nClear group structure.\nCommunity-rooted operating perspective.\nCommercial language grounded in execution rather than hype.",
        accent_color=ACCENT,
    )

    # Slide 3 why story philosophy
    slide = prs.slides.add_slide(blank)
    add_slide_base(slide, dark=True)
    add_ppt_title(slide, "Why We Exist", "Purpose before products. Trust in delivery. Growth with discipline.", dark=True)
    add_card(
        slide,
        0.7,
        1.7,
        3.9,
        1.65,
        "Our why",
        "Communities deserve quality, professionalism and dignity in the services they rely on every day.",
        accent_color=ACCENT,
        dark=True,
    )
    add_card(
        slide,
        0.7,
        3.6,
        5.1,
        2.1,
        "Our story",
        "Founder and CEO Masego Mafoko built Kgolaentle Holdings on the belief that communities in and around Rustenburg deserved the same standard of service found anywhere else in the country.",
        accent_color=PRIMARY,
        dark=True,
    )
    add_card(
        slide,
        6.05,
        1.7,
        6.15,
        4.0,
        "Our philosophy",
        "Strong businesses are built from the inside out. Purpose comes first. Trust is earned in delivery. Growth must be disciplined. Excellence is the result of consistent standards applied over time.",
        accent_color=ACCENT,
        dark=True,
    )

    # Slide 4 what sets apart
    slide = prs.slides.add_slide(blank)
    add_slide_base(slide, dark=False)
    add_ppt_title(slide, "What Sets The Group Apart", "A holdings model with clear standards and long-term intent")
    add_ppt_bullets(
        slide,
        [
            "Disciplined diversification rather than opportunistic expansion",
            "Local grounding that keeps decision-making close to real community needs",
            "Service shaped by trust, communication and dependable execution",
            "Commercial maturity without inflated claims or empty corporate language",
        ],
        left=0.8,
        top=1.7,
        width=5.8,
        height=2.4,
        size=12,
    )
    add_card(
        slide,
        7.0,
        1.5,
        5.5,
        2.0,
        "Group operating model",
        "Each portfolio has its own customer need, operating rhythm and commercial role. What connects them is shared leadership direction, high service expectations and a long-term view of growth.",
        accent_color=PRIMARY,
    )
    add_card(slide, 7.0, 3.9, 2.55, 1.5, "Trust", "Reliability, transparency and follow-through.", accent_color=ACCENT)
    add_card(slide, 9.75, 3.9, 2.75, 1.5, "Collaboration", "Better outcomes through strong partnership and teamwork.", accent_color=PRIMARY)
    add_card(slide, 7.0, 5.65, 5.5, 1.2, "Continuous improvement", "A commitment to keep refining service, systems and execution quality.", accent_color=ACCENT)

    # Portfolio slides
    for portfolio in PORTFOLIOS:
        slide = prs.slides.add_slide(blank)
        add_slide_base(slide, dark=False)
        add_ppt_title(slide, portfolio["title"], portfolio["tagline"])
        slide.shapes.add_picture(str(portfolio["image"]), PptInches(7.55), PptInches(1.2), width=PptInches(5.0), height=PptInches(4.8))
        frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PptInches(0.75), PptInches(1.55), PptInches(6.2), PptInches(4.95))
        frame.fill.solid()
        frame.fill.fore_color.rgb = ppt_rgb(WARM_GRAY)
        frame.line.color.rgb = ppt_rgb(PRIMARY)
        frame.line.width = PptPt(1.6)
        tf = frame.text_frame
        tf.clear()
        for idx, block in enumerate(
            [
                ("What it does", portfolio["summary"]),
                ("Who it serves", portfolio["serves"]),
                ("Value delivered", portfolio["value"]),
            ]
        ):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            r = p.add_run()
            r.text = block[0]
            r.font.name = "Georgia"
            r.font.size = PptPt(13)
            r.font.bold = True
            r.font.color.rgb = ppt_rgb(SECONDARY)
            p.space_after = PptPt(2)
            p2 = tf.add_paragraph()
            r2 = p2.add_run()
            r2.text = block[1]
            r2.font.name = "Arial"
            r2.font.size = PptPt(10.5)
            r2.font.color.rgb = ppt_rgb(BODY)
            p2.space_after = PptPt(8)
        add_ppt_bullets(slide, portfolio["bullets"], left=7.7, top=6.1, width=4.75, height=1.0, size=9.5)

    # Leadership and close
    slide = prs.slides.add_slide(blank)
    add_slide_base(slide, dark=True)
    add_ppt_title(slide, "Founder-Led Leadership", "Conviction, care, discipline and long-term intent", dark=True)
    add_card(
        slide,
        0.8,
        1.7,
        6.2,
        3.45,
        "Masego Mafoko | Founder and CEO",
        "With over ten years of experience in communications, public relations and marketing, Masego Mafoko has shaped Kgolaentle Holdings around a simple standard: communities deserve service that is professional, reliable and dignified. Her leadership has created a group with a clear why and a disciplined approach to growth.",
        accent_color=ACCENT,
        dark=True,
    )
    add_card(
        slide,
        7.35,
        1.7,
        5.1,
        2.0,
        "Commitment to community and excellence",
        "Commercial success and community impact should strengthen each other. When the business shows up with care and keeps its word it builds confidence, dignity and long-term trust.",
        accent_color=PRIMARY,
        dark=True,
    )
    add_card(
        slide,
        7.35,
        4.0,
        5.1,
        2.05,
        "Contact",
        "Blairgowrie Section, Chaneng\nNorth West, 0310, South Africa\n+27 (0) 87 093 7316\ninfo@kgolaentle.com\nwww.kgolaentle.com",
        accent_color=ACCENT,
        dark=True,
    )

    prs.save(PPTX_PATH)


if __name__ == "__main__":
    build_docx()
    build_pptx()
    print(f"Generated {DOCX_PATH}")
    print(f"Generated {PPTX_PATH}")
