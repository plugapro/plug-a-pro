"""Build the Plug A Pro BNI 60-second visitor-intro slide (16:9 .pptx).

Layout (left → right):
  - NAVY content panel (~58%): logo, headline, sub-line, short description,
    three quick value bullets, BNI footer chip.
  - PHOTO panel (~42%): full-bleed provider hero, with a floating white QR
    card overlaid at the bottom-right carrying the QR + WhatsApp + URL.

This deliberately mirrors the West Rand pilot flyer's hierarchy but is
re-laid out for landscape projection at conference scale: large headline,
single visual focal point, one obvious CTA (the QR card).

Brand palette (sampled from Plug A Pro flyers):
  navy        #0B1B3F
  navy-soft   #112C58
  green       #1E8E5A
  green-soft  #2BB673
  accent      #E07A2C
  text        #1F2A44
  muted-text  #5E6A85
  paper       #FFFFFF
"""
from __future__ import annotations

from pathlib import Path

from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

ROOT = Path("/Users/shimane/Projects/Plug A Pro/docs/INF Presentation/build")
OUT = Path(
    "/Users/shimane/Projects/Plug A Pro/docs/INF Presentation/plugapro_bni_intro_slide.pptx"
)

NAVY = RGBColor(0x0B, 0x1B, 0x3F)
NAVY_SOFT = RGBColor(0x11, 0x2C, 0x58)
NAVY_DEEP = RGBColor(0x06, 0x11, 0x2A)
GREEN = RGBColor(0x1E, 0x8E, 0x5A)
GREEN_SOFT = RGBColor(0x2B, 0xB6, 0x73)
ACCENT = RGBColor(0xE0, 0x7A, 0x2C)
TEXT = RGBColor(0x1F, 0x2A, 0x44)
MUTED = RGBColor(0x6B, 0x76, 0x90)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
PAPER_SOFT = RGBColor(0xF6, 0xF8, 0xFC)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Image prep
# ---------------------------------------------------------------------------


def prep_hero_panel() -> Path:
    """Make a tall landscape-cropped provider hero that fills the photo panel.

    Source image is wide (~16:9) with the man on the right and a softly faded
    warehouse on the left. We crop a slightly tighter window around the man,
    then add a soft navy gradient on the right edge so the floating QR card
    can sit on it without colour clashes.
    """
    src = Image.open(ROOT / "hero_provider_with_sign.png").convert("RGB")
    sw, sh = src.size
    target_w, target_h = 1800, 2160  # close to 5:6 portrait
    # Cover-fit using full source
    scale = max(target_w / sw, target_h / sh)
    new_size = (int(sw * scale), int(sh * scale))
    resized = src.resize(new_size, Image.LANCZOS)
    # Horizontal bias chosen so the HANDYMAN / PAINTER / TILER sign AND the
    # provider both stay in frame. ~0.62 favours the right (where the man is)
    # while keeping the sign's left edge readable.
    horiz_overflow = max(0, resized.size[0] - target_w)
    left = int(horiz_overflow * 0.62)
    top = (resized.size[1] - target_h) // 2
    final = resized.crop((left, top, left + target_w, top + target_h))

    # Subtle bottom vignette so the QR card grounds visually
    vignette = Image.new("RGBA", (target_w, target_h), (0, 0, 0, 0))
    vd = ImageDraw.Draw(vignette)
    for i in range(int(target_h * 0.35)):
        y = target_h - 1 - i
        alpha = int(110 * (i / (target_h * 0.35)))
        vd.line([(0, y), (target_w, y)], fill=(8, 18, 42, alpha))
    final_rgba = final.convert("RGBA")
    final_rgba = Image.alpha_composite(final_rgba, vignette)

    # Subtle left-edge fade into navy so the seam against the content panel
    # is butter-smooth rather than a hard cut.
    edge = Image.new("RGBA", (target_w, target_h), (0, 0, 0, 0))
    ed = ImageDraw.Draw(edge)
    fade_w = int(target_w * 0.18)
    for x in range(fade_w):
        alpha = int(230 * (1 - x / fade_w) ** 1.4)
        ed.line([(x, 0), (x, target_h)], fill=(11, 27, 63, alpha))
    final_rgba = Image.alpha_composite(final_rgba, edge)

    out = ROOT / "hero_panel.png"
    final_rgba.convert("RGB").save(out, "PNG", quality=92)
    return out


def prep_navy_panel() -> Path:
    """Render the navy left-content background with a subtle radial highlight
    so it doesn't read as flat-fill in projection."""
    w, h = 1800, 2160
    img = Image.new("RGB", (w, h), (11, 27, 63))
    # Vertical gradient
    grad = Image.new("RGB", (1, h))
    for y in range(h):
        t = y / h
        r = int(11 + (6 - 11) * t)
        g = int(27 + (17 - 27) * t)
        b = int(63 + (42 - 63) * t)
        grad.putpixel((0, y), (r, g, b))
    img.paste(grad.resize((w, h)), (0, 0))

    # Soft highlight in upper-left for depth
    highlight = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    hd = ImageDraw.Draw(highlight)
    cx, cy, rad = int(w * 0.20), int(h * 0.18), int(w * 0.55)
    for r in range(rad, 0, -2):
        a = int(38 * (1 - r / rad))
        hd.ellipse([cx - r, cy - r, cx + r, cy + r], fill=(255, 255, 255, a))
    highlight = highlight.filter(ImageFilter.GaussianBlur(radius=50))
    img.paste(highlight, (0, 0), highlight)

    # Diagonal green accent slash bottom-left
    slash = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    sd = ImageDraw.Draw(slash)
    sd.polygon(
        [(0, h), (0, int(h * 0.88)), (int(w * 0.42), h)],
        fill=(43, 182, 115, 60),
    )
    slash = slash.filter(ImageFilter.GaussianBlur(radius=30))
    img.paste(slash, (0, 0), slash)

    out = ROOT / "navy_panel.png"
    img.save(out, "PNG", quality=92)
    return out


# ---------------------------------------------------------------------------
# pptx helpers
# ---------------------------------------------------------------------------


def add_rect(slide, x, y, w, h, fill: RGBColor, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_rounded(slide, x, y, w, h, fill: RGBColor, line=None, line_w=None, corner=0.06):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = corner
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        if line_w is not None:
            shp.line.width = line_w
    shp.shadow.inherit = False
    return shp


def add_text(
    slide,
    x,
    y,
    w,
    h,
    runs,
    *,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor

    for i, run_def in enumerate(runs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if "space_before" in run_def:
            p.space_before = Pt(run_def["space_before"])
        if "space_after" in run_def:
            p.space_after = Pt(run_def["space_after"])
        if "line_spacing" in run_def:
            p.line_spacing = run_def["line_spacing"]
        r = p.add_run()
        r.text = run_def["text"]
        f = r.font
        f.size = Pt(run_def.get("size", 14))
        f.bold = run_def.get("bold", False)
        f.italic = run_def.get("italic", False)
        f.name = run_def.get("font", "Calibri")
        f.color.rgb = run_def.get("color", TEXT)
    return tb


# ---------------------------------------------------------------------------
# Slide composition
# ---------------------------------------------------------------------------


def build() -> None:
    navy_panel = prep_navy_panel()
    hero_panel = prep_hero_panel()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # ---- 1. Backgrounds ----
    content_w = Inches(7.7)
    photo_x = content_w
    photo_w = SLIDE_W - content_w

    slide.shapes.add_picture(str(navy_panel), Inches(0), Inches(0), content_w, SLIDE_H)
    slide.shapes.add_picture(str(hero_panel), photo_x, Inches(0), photo_w, SLIDE_H)

    # Thin green seam between zones
    add_rect(slide, content_w - Inches(0.025), Inches(0), Inches(0.05), SLIDE_H, GREEN_SOFT)

    # ---- 2. Logo + brand strip top-left ----
    slide.shapes.add_picture(
        str(ROOT / "logo.png"),
        Inches(0.55),
        Inches(0.50),
        height=Inches(0.78),
    )

    # Small tagline under logo
    add_text(
        slide,
        Inches(0.58),
        Inches(1.32),
        Inches(6.8),
        Inches(0.30),
        [
            {
                "text": "LOCAL SERVICES MARKETPLACE  ·  WHATSAPP-FIRST",
                "size": 11,
                "bold": True,
                "color": GREEN_SOFT,
                "font": "Calibri",
            }
        ],
    )

    # ---- 3. Headline ----
    add_text(
        slide,
        Inches(0.55),
        Inches(1.95),
        Inches(7.0),
        Inches(2.4),
        [
            {
                "text": "Connecting Local Skills",
                "size": 56,
                "bold": True,
                "color": PAPER,
                "font": "Calibri",
                "line_spacing": 1.0,
            },
            {
                "text": "to Local Work.",
                "size": 56,
                "bold": True,
                "color": GREEN_SOFT,
                "font": "Calibri",
                "line_spacing": 1.0,
            },
        ],
    )

    # ---- 4. Pill ----
    pill_y = Inches(4.30)
    pill_w = Inches(4.6)
    add_rounded(
        slide,
        Inches(0.55),
        pill_y,
        pill_w,
        Inches(0.55),
        NAVY_DEEP,
        line=GREEN_SOFT,
        line_w=Pt(1.25),
        corner=0.5,
    )
    add_text(
        slide,
        Inches(0.55),
        pill_y,
        pill_w,
        Inches(0.55),
        [
            {
                "text": "Home & business maintenance, booked on WhatsApp",
                "size": 14,
                "bold": True,
                "color": PAPER,
                "font": "Calibri",
            }
        ],
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # ---- 5. Description ----
    add_text(
        slide,
        Inches(0.55),
        Inches(5.10),
        Inches(7.0),
        Inches(1.4),
        [
            {
                "text": "Plug A Pro is a simple way for households and "
                "businesses to find local service providers — from plumbing "
                "and electrical work to handywork, painting, and beyond. "
                "Customers send 'JOIN' on WhatsApp, we route the job to "
                "nearby providers, and updates flow back as the work moves.",
                "size": 14,
                "color": RGBColor(0xDD, 0xE3, 0xF1),
                "font": "Calibri",
                "line_spacing": 1.30,
            }
        ],
    )

    # ---- 6. BNI footer chip ----
    bni_y = Inches(6.78)
    add_rounded(
        slide,
        Inches(0.55),
        bni_y,
        Inches(4.2),
        Inches(0.42),
        ACCENT,
        corner=0.5,
    )
    add_text(
        slide,
        Inches(0.55),
        bni_y,
        Inches(4.2),
        Inches(0.42),
        [
            {
                "text": "BNI Prosper-Us  ·  Visitor Introduction  ·  Lebogang",
                "size": 11,
                "bold": True,
                "color": PAPER,
                "font": "Calibri",
            }
        ],
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # ---- 7. Floating QR / contact card on the right photo panel ----
    card_w = Inches(2.95)
    card_h = Inches(3.65)
    card_x = SLIDE_W - card_w - Inches(0.25)
    card_y = SLIDE_H - card_h - Inches(0.25)

    # Drop-shadow approximation: a slightly offset dark rectangle behind the card
    add_rounded(
        slide,
        card_x + Inches(0.08),
        card_y + Inches(0.10),
        card_w,
        card_h,
        RGBColor(0x05, 0x0D, 0x22),
        corner=0.04,
    )
    add_rounded(
        slide,
        card_x,
        card_y,
        card_w,
        card_h,
        PAPER,
        line=NAVY,
        line_w=Pt(1.25),
        corner=0.04,
    )

    # Card header — green band
    header_h = Inches(0.55)
    add_rounded(
        slide,
        card_x,
        card_y,
        card_w,
        header_h,
        GREEN,
        corner=0.04,
    )
    # Hide the bottom-rounded sides of the header by masking with a square
    add_rect(
        slide,
        card_x,
        card_y + Inches(0.27),
        card_w,
        Inches(0.30),
        GREEN,
    )
    add_text(
        slide,
        card_x,
        card_y,
        card_w,
        header_h,
        [
            {
                "text": "SCAN TO LEARN MORE",
                "size": 14,
                "bold": True,
                "color": PAPER,
                "font": "Calibri",
            }
        ],
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # QR
    qr_size = Inches(1.85)
    qr_x = card_x + (card_w - qr_size) / 2
    qr_y = card_y + Inches(0.65)
    slide.shapes.add_picture(
        str(ROOT / "qr_plugapro.png"),
        qr_x,
        qr_y,
        width=qr_size,
        height=qr_size,
    )

    # Below QR — bigger "or join the local rollout"
    add_text(
        slide,
        card_x + Inches(0.30),
        qr_y + qr_size + Inches(0.05),
        card_w - Inches(0.60),
        Inches(0.38),
        [
            {
                "text": "or join the local rollout",
                "size": 13,
                "italic": True,
                "color": MUTED,
                "font": "Calibri",
            }
        ],
        align=PP_ALIGN.CENTER,
    )

    # Divider
    add_rect(
        slide,
        card_x + Inches(0.30),
        qr_y + qr_size + Inches(0.45),
        card_w - Inches(0.60),
        Emu(8000),
        RGBColor(0xE0, 0xE4, 0xEE),
    )

    # WhatsApp row
    contact_y = qr_y + qr_size + Inches(0.55)
    add_rounded(
        slide,
        card_x + Inches(0.30),
        contact_y,
        Inches(0.38),
        Inches(0.38),
        GREEN,
        corner=0.5,
    )
    add_text(
        slide,
        card_x + Inches(0.30),
        contact_y,
        Inches(0.38),
        Inches(0.38),
        [
            {
                "text": "WA",
                "size": 10,
                "bold": True,
                "color": PAPER,
                "font": "Calibri",
            }
        ],
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        card_x + Inches(0.80),
        contact_y - Inches(0.04),
        card_w - Inches(1.10),
        Inches(0.50),
        [
            {
                "text": "+27 69 355 2447",
                "size": 18,
                "bold": True,
                "color": NAVY,
                "font": "Calibri",
            }
        ],
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # Website row
    web_y = contact_y + Inches(0.50)
    add_rounded(
        slide,
        card_x + Inches(0.30),
        web_y,
        Inches(0.38),
        Inches(0.38),
        NAVY,
        corner=0.5,
    )
    add_text(
        slide,
        card_x + Inches(0.30),
        web_y,
        Inches(0.38),
        Inches(0.38),
        [
            {
                "text": ".za",
                "size": 9,
                "bold": True,
                "color": PAPER,
                "font": "Calibri",
            }
        ],
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        card_x + Inches(0.80),
        web_y - Inches(0.04),
        card_w - Inches(1.10),
        Inches(0.50),
        [
            {
                "text": "plugapro.co.za",
                "size": 18,
                "bold": True,
                "color": NAVY,
                "font": "Calibri",
            }
        ],
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # ---- 8. Speaker notes ----
    notes = slide.notes_slide.notes_text_frame
    notes.text = (
        "60-second BNI visitor intro — Lebogang.\n\n"
        "Plug A Pro is a WhatsApp-first local services marketplace. "
        "We connect customers with local service providers for home and "
        "business maintenance — plumbing, electrical, handywork, painting, "
        "and more.\n\n"
        "Customers send 'JOIN' on WhatsApp or scan the QR code on the slide. "
        "We help them describe the job, route it to nearby providers, and "
        "keep them updated as the work moves.\n\n"
        "Good referral for me today: anyone with a maintenance job they keep "
        "putting off, or any local service provider who wants more "
        "consistent work. Scan to learn more or join the local rollout."
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
